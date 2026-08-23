"""PowerPoint builder: widescreen grid, RTL, sections, click/hover zoom."""

from __future__ import annotations

from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

from .image_processor import compress_image_to_bytes, get_image_size
from .models import BuildSettings
from .scanner import ImageGroup, PresentationJob

TITLE_COLOR = RGBColor(0, 0, 0)
MUTED_COLOR = RGBColor(80, 80, 80)
ACCENT_COLOR = RGBColor(15, 61, 46)
BORDER_COLOR = RGBColor(180, 180, 180)

MARGIN_X = Inches(0.35)
GUTTER = Inches(0.25)
TITLE_BAND = Inches(0.35)
LOGO_SIZE = Inches(0.85)

HEADER_TOP = Inches(0.0)
HEADER_BOTTOM = Inches(1.2)
GRID_ROW1_TOP = Inches(1.2)
GRID_ROW1_BOTTOM = Inches(4.0)
GRID_ROW2_TOP = Inches(4.0)
GRID_ROW2_BOTTOM = Inches(6.8)
FOOTER_TOP = Inches(6.8)
FOOTER_BOTTOM = Inches(7.5)


from app.i18n import set_build_slide_language, t_slide


def set_rtl(paragraph, *, rtl: bool = True) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    if rtl:
        pPr.set("rtl", "1")
    elif pPr.get("rtl") is not None:
        del pPr.attrib["rtl"]


def _paragraph_rtl(settings: BuildSettings) -> bool:
    return settings.slide_language != "en"


def _style_run(run, settings: BuildSettings, *, size_pt: float, bold: bool = False, color: RGBColor = TITLE_COLOR) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = settings.font_name
    rPr = run._r.get_or_add_rPr()
    existing_cs = rPr.find(qn("a:cs"))
    if existing_cs is not None:
        rPr.remove(existing_cs)
    cs = OxmlElement("a:cs")
    cs.set("typeface", settings.font_name)
    rPr.append(cs)


def _add_textbox(slide, left, top, width, height, text: str, settings: BuildSettings, *, size_pt: float, bold: bool = False, color: RGBColor = TITLE_COLOR, align=None):
    rtl = _paragraph_rtl(settings)
    if align is None:
        align = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_rtl(p, rtl=rtl)
    run = p.add_run()
    run.text = text
    _style_run(run, settings, size_pt=size_pt, bold=bold, color=color)
    return shape


def _fit_image_in_box(img_w, img_h, box_left, box_top, box_width, box_height) -> tuple:
    box_w = float(box_width)
    box_h = float(box_height)
    if img_w <= 0 or img_h <= 0:
        return box_left, box_top, box_width, box_height
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio >= box_ratio:
        final_w, final_h = box_w, box_w / img_ratio
    else:
        final_h, final_w = box_h, box_h * img_ratio
    left = float(box_left) + (box_w - final_w) / 2
    top = float(box_top) + (box_h - final_h) / 2
    return Emu(int(left)), Emu(int(top)), Emu(int(final_w)), Emu(int(final_h))


def _cell_geometry(index: int, slide_width, images_per_slide: int) -> tuple:
    slide_w = float(slide_width)
    margin = float(MARGIN_X)
    gutter = float(GUTTER)
    cols = 2 if images_per_slide > 1 else 1
    rows = max(1, (images_per_slide + cols - 1) // cols)
    usable_w = slide_w - 2 * margin - gutter * (cols - 1)
    cell_w = usable_w / cols

    col = index % cols
    row = index // cols

    body_top = float(GRID_ROW1_TOP)
    body_bottom = float(GRID_ROW2_BOTTOM if rows > 1 else GRID_ROW1_BOTTOM)
    body_h = body_bottom - body_top
    row_h = body_h / rows

    left = margin + col * (cell_w + gutter)
    band_top = body_top + row * row_h
    title_h = float(TITLE_BAND)
    title_box = (Emu(int(left)), Emu(int(band_top)), Emu(int(cell_w)), Emu(int(title_h)))
    img_top = band_top + title_h
    img_h = band_top + row_h - img_top
    image_box = (Emu(int(left)), Emu(int(img_top)), Emu(int(cell_w)), Emu(int(img_h)))
    return title_box, image_box


def _decorate_picture(pic, settings: BuildSettings) -> None:
    if settings.enable_image_border:
        pic.line.color.rgb = BORDER_COLOR
        pic.line.width = Pt(0.75)
    if settings.enable_image_shadow:
        pic.shadow.inherit = False


def _link_shape_to_slide(shape, target_slide) -> None:
    shape.click_action.target_slide = target_slide


def _link_shape_hover(shape, target_slide) -> None:
    """Hover jump via OpenXML (avoids python-pptx hover ActionSetting bug)."""
    cNvPr = shape._element._nvXxPr.cNvPr
    hlink = cNvPr.get_or_add_hlinkHover()
    hlink.action = "ppaction://hlinksldjump"
    hlink.rId = shape.part.relate_to(target_slide.part, RT.SLIDE)


def _compress(image_path: Path, settings: BuildSettings):
    return compress_image_to_bytes(
        image_path,
        max_dimension=settings.max_dimension,
        jpeg_quality=settings.jpeg_quality,
    )


def _add_logo(slide, path: Path | None, left, top, settings: BuildSettings, size=LOGO_SIZE) -> None:
    if path is None or not Path(path).is_file():
        return
    try:
        buffer = _compress(path, settings)
        w_px, h_px = get_image_size(path)
        left_f, top_f, w_f, h_f = _fit_image_in_box(w_px, h_px, left, top, size, size)
        slide.shapes.add_picture(buffer, left_f, top_f, w_f, h_f)
    except Exception:
        pass


def _header_title_box(slide_width):
    title_left = Emu(int(float(MARGIN_X) + float(LOGO_SIZE) + float(GUTTER)))
    title_width = Emu(int(float(slide_width) - 2 * (float(MARGIN_X) + float(LOGO_SIZE) + float(GUTTER))))
    return title_left, title_width


def _draw_header(slide, title: str, settings: BuildSettings, logo_right, logo_left, slide_width, *, subtitle: str | None = None) -> None:
    header_h = float(HEADER_BOTTOM) - float(HEADER_TOP)
    logo_top = Emu(int(float(HEADER_TOP) + (header_h - float(LOGO_SIZE)) / 2))
    _add_logo(slide, logo_right, Emu(int(float(slide_width) - float(MARGIN_X) - float(LOGO_SIZE))), logo_top, settings)
    _add_logo(slide, logo_left, MARGIN_X, logo_top, settings)
    title_left, title_width = _header_title_box(slide_width)
    if subtitle:
        _add_textbox(slide, title_left, Emu(int(float(HEADER_TOP) + 0.12 * 914400)), title_width, Emu(int(0.45 * 914400)), title, settings, size_pt=settings.title_font_size - 4, bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(slide, title_left, Emu(int(float(HEADER_TOP) + 0.55 * 914400)), title_width, Emu(int(0.4 * 914400)), subtitle, settings, size_pt=13, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
    else:
        _add_textbox(slide, title_left, Emu(int(float(HEADER_TOP) + 0.25 * 914400)), title_width, Emu(int(0.7 * 914400)), title, settings, size_pt=settings.title_font_size, bold=True, align=PP_ALIGN.CENTER)


def _draw_footer(slide, footer_text: str, settings: BuildSettings, slide_width) -> None:
    if not footer_text.strip():
        return
    footer_h = float(FOOTER_BOTTOM) - float(FOOTER_TOP)
    _add_textbox(slide, MARGIN_X, FOOTER_TOP, Emu(int(float(slide_width) - 2 * float(MARGIN_X))), Emu(int(footer_h)), footer_text, settings, size_pt=settings.footer_font_size, color=MUTED_COLOR, align=PP_ALIGN.CENTER)


def _caption_for(image_path: Path, settings: BuildSettings) -> str:
    return image_path.stem if settings.caption_from_filename else ""


def _add_detail_slide(prs, blank, image_path: Path, settings: BuildSettings, *, person_title: str, footer_text: str, logo_right, logo_left, slide_width, slide_height, return_slide=None):
    slide = prs.slides.add_slide(blank)
    caption = _caption_for(image_path, settings)
    _draw_header(slide, person_title, settings, logo_right, logo_left, slide_width)
    _draw_footer(slide, footer_text, settings, slide_width)

    margin = float(MARGIN_X)
    top = float(HEADER_BOTTOM) + margin * 0.5
    bottom = float(FOOTER_TOP) - margin
    caption_band = 0.45 * 914400 if caption else 0
    box_left = margin
    box_w = float(slide_width) - 2 * margin
    box_h = bottom - top - caption_band

    buffer = _compress(image_path, settings)
    w_px, h_px = get_image_size(image_path)
    left, top_e, width, height = _fit_image_in_box(w_px, h_px, Emu(int(box_left)), Emu(int(top)), Emu(int(box_w)), Emu(int(box_h)))
    pic = slide.shapes.add_picture(buffer, left, top_e, width, height)
    _decorate_picture(pic, settings)

    if caption:
        cap_top = top + box_h + 0.08 * 914400
        _add_textbox(
            slide,
            Emu(int(box_left)),
            Emu(int(cap_top)),
            Emu(int(box_w)),
            Emu(int(caption_band)),
            caption,
            settings,
            size_pt=settings.caption_font_size,
            color=MUTED_COLOR,
            align=PP_ALIGN.CENTER,
        )

    if return_slide is not None:
        _link_shape_to_slide(pic, return_slide)
    return slide


def _add_image_cell(slide, prs, blank, image_path: Path, cell_index: int, settings: BuildSettings, *, person_title: str, footer_text: str, logo_right, logo_left, slide_width, slide_height) -> None:
    title_box, image_box = _cell_geometry(cell_index, slide_width, settings.images_per_slide)
    caption = _caption_for(image_path, settings)
    if caption:
        _add_textbox(slide, *title_box, caption, settings, size_pt=settings.caption_font_size, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    buffer = _compress(image_path, settings)
    w_px, h_px = get_image_size(image_path)
    left, top, width, height = _fit_image_in_box(w_px, h_px, *image_box)
    pic = slide.shapes.add_picture(buffer, left, top, width, height)
    _decorate_picture(pic, settings)

    if settings.enable_image_zoom:
        detail = _add_detail_slide(
            prs,
            blank,
            image_path,
            settings,
            person_title=person_title,
            footer_text=footer_text,
            logo_right=logo_right,
            logo_left=logo_left,
            slide_width=slide_width,
            slide_height=slide_height,
            return_slide=slide,
        )
        _link_shape_to_slide(pic, detail)
        if settings.enable_image_zoom and settings.enable_hover_zoom:
            _link_shape_hover(pic, detail)


def _add_section_divider(prs, blank, settings, *, person_title, section_name, section_index, section_total, footer_text, logo_right, logo_left, slide_width):
    slide = prs.slides.add_slide(blank)
    _draw_header(slide, person_title, settings, logo_right, logo_left, slide_width)
    _draw_footer(slide, footer_text, settings, slide_width)
    _add_textbox(
        slide,
        MARGIN_X,
        Emu(int(2.6 * 914400)),
        Emu(int(float(slide_width) - 2 * float(MARGIN_X))),
        Emu(int(0.5 * 914400)),
        t_slide("pptx.section.n_of_m", n=section_index, m=section_total),
        settings,
        size_pt=14,
        color=MUTED_COLOR,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(slide, MARGIN_X, Emu(int(3.2 * 914400)), Emu(int(float(slide_width) - 2 * float(MARGIN_X))), Emu(int(1.2 * 914400)), section_name, settings, size_pt=32, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)


def _add_grid_slides(prs, blank, images, settings, *, title, subtitle, footer_text, logo_right, logo_left, slide_width, slide_height, should_cancel):
    per = max(1, min(4, settings.images_per_slide))
    for i in range(0, len(images), per):
        if should_cancel and should_cancel():
            raise InterruptedError(t_slide("pptx.err.cancelled"))
        chunk = images[i : i + per]
        slide = prs.slides.add_slide(blank)
        _draw_header(slide, title, settings, logo_right, logo_left, slide_width, subtitle=subtitle)
        _draw_footer(slide, footer_text, settings, slide_width)
        for cell_index, img_path in enumerate(chunk):
            if should_cancel and should_cancel():
                raise InterruptedError(t_slide("pptx.err.cancelled"))
            try:
                _add_image_cell(
                    slide,
                    prs,
                    blank,
                    img_path,
                    cell_index,
                    settings,
                    person_title=title,
                    footer_text=footer_text,
                    logo_right=logo_right,
                    logo_left=logo_left,
                    slide_width=slide_width,
                    slide_height=slide_height,
                )
            except Exception as exc:
                raise RuntimeError(t_slide("pptx.err.image", name=img_path.name, exc=exc)) from exc


def build_presentation_from_job(
    job: PresentationJob,
    output_path: Path,
    *,
    settings: BuildSettings | None = None,
    should_cancel: Callable[[], bool] | None = None,
) -> Path:
    cfg = settings or BuildSettings()
    set_build_slide_language(cfg.slide_language)
    prs = Presentation()
    slide_width = Inches(cfg.slide_width_inches)
    slide_height = Inches(cfg.slide_height_inches)
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    blank = prs.slide_layouts[6]

    use_sections = cfg.enable_section_dividers and job.grouped and len(job.groups) > 1
    section_total = len(job.groups)

    for section_index, group in enumerate(job.groups, start=1):
        if should_cancel and should_cancel():
            raise InterruptedError(t_slide("pptx.err.cancelled"))
        if use_sections:
            _add_section_divider(
                prs,
                blank,
                cfg,
                person_title=job.name,
                section_name=group.name,
                section_index=section_index,
                section_total=section_total,
                footer_text=cfg.footer_text,
                logo_right=cfg.logo_right,
                logo_left=cfg.logo_left,
                slide_width=slide_width,
            )
            subtitle = group.name
        else:
            subtitle = None
        _add_grid_slides(
            prs,
            blank,
            group.images,
            cfg,
            title=job.name,
            subtitle=subtitle,
            footer_text=cfg.footer_text,
            logo_right=cfg.logo_right,
            logo_left=cfg.logo_left,
            slide_width=slide_width,
            slide_height=slide_height,
            should_cancel=should_cancel,
        )

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def build_presentation(images, output_path, *, title: str, settings: BuildSettings | None = None, should_cancel=None):
    job = PresentationJob(name=title, source=Path(output_path).parent, groups=[ImageGroup(name=title, images=list(images))], grouped=False)
    cfg = settings or BuildSettings()
    return build_presentation_from_job(job, output_path, settings=cfg, should_cancel=should_cancel)
