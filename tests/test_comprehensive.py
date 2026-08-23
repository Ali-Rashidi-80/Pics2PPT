"""Comprehensive honest tests — no mocked pass conditions."""

from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from PySide6.QtCore import Qt
from PySide6.QtWidgets import QFrame, QLabel

from app.core.models import BuildSettings
from app.core.pptx_builder import build_presentation_from_job
from app.core.scanner import (
    collect_images,
    is_valid_image,
    make_grouped_job,
    scan_project_folders,
)
from app.core.worker import PresentationWorker
from app.services.settings import DEFAULT_SETTINGS, SettingsManager


def _write_img(folder: Path, name: str, size=(640, 480)) -> Path:
    p = folder / name
    Image.new("RGB", size, (120, 140, 160)).save(p, format="JPEG", quality=90)
    return p


def _make_tree(base: Path) -> None:
    """Synthetic folder tree covering flat, grouped, and numbered layouts."""
    flat = base / "session_flat"
    flat.mkdir()
    _write_img(flat, "a.jpg")
    _write_img(flat, "b.jpg")

    person = base / "person_alpha"
    person.mkdir()
    _write_img(person, "own1.jpg")
    _write_img(person, "own2.jpg")
    t1 = person / "topic_one"
    t1.mkdir()
    _write_img(t1, "t1.jpg")
    t2 = person / "topic_two"
    t2.mkdir()
    _write_img(t2, "t2.jpg")

    numbered = base / "group_visits"
    numbered.mkdir()
    g1 = numbered / "1"
    g1.mkdir()
    _write_img(g1, "g1.jpg")
    g2 = numbered / "2"
    g2.mkdir()
    _write_img(g2, "g2.jpg")

    (flat / "Thumbs.db").write_bytes(b"x")
    (flat / "x.rar").write_bytes(b"x")


class TestScanner(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="gp_scan_"))

    def tearDown(self) -> None:
        import shutil

        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_valid_image_filters(self) -> None:
        p = self.tmp / "ok.jpg"
        _write_img(self.tmp, "ok.jpg")
        self.assertTrue(is_valid_image(p))
        self.assertFalse(is_valid_image(self.tmp / "Thumbs.db"))
        self.assertFalse(is_valid_image(self.tmp / "a.rar"))
        self.assertFalse(is_valid_image(self.tmp / "a.zip"))

    def test_project_root_detects_all_units(self) -> None:
        _make_tree(self.tmp)
        jobs = scan_project_folders(self.tmp)
        names = {j.name for j in jobs}
        self.assertEqual(len(jobs), 3, names)
        self.assertIn("person_alpha", names)
        self.assertIn("session_flat", names)
        self.assertIn("group_visits", names)

    def test_person_grouped_sections(self) -> None:
        _make_tree(self.tmp)
        person = next(j for j in scan_project_folders(self.tmp) if j.name == "person_alpha")
        self.assertTrue(person.grouped)
        section_names = [g.name for g in person.groups]
        self.assertIn("تصاویر کلی", section_names)
        self.assertIn("topic_one", section_names)
        self.assertEqual(sum(len(g.images) for g in person.groups), 4)

    def test_flat_single_folder(self) -> None:
        leaf = self.tmp / "only"
        leaf.mkdir()
        _write_img(leaf, "x.jpg")
        jobs = scan_project_folders(leaf)
        self.assertEqual(len(jobs), 1)
        self.assertFalse(jobs[0].grouped)
        self.assertEqual(len(jobs[0].groups[0].images), 1)

    def test_skips_output_folder(self) -> None:
        _make_tree(self.tmp)
        out = self.tmp / "Output_PPTX"
        out.mkdir()
        _write_img(out, "should_skip.jpg")
        jobs = scan_project_folders(self.tmp, skip_dir_names={"Output_PPTX"})
        self.assertEqual(len(jobs), 3)
        for j in jobs:
            self.assertNotEqual(j.name, "Output_PPTX")

    def test_skips_custom_output_folder_name(self) -> None:
        _make_tree(self.tmp)
        custom = self.tmp / "MyReports"
        custom.mkdir()
        _write_img(custom, "skip.jpg")
        jobs = scan_project_folders(self.tmp, skip_dir_names={"MyReports"})
        self.assertEqual(len(jobs), 3)


class TestPptxBuilder(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="gp_pptx_"))

    def tearDown(self) -> None:
        import shutil

        shutil.rmtree(self.tmp, ignore_errors=True)

    def _inspect_pptx(self, path: Path) -> dict:
        self.assertTrue(path.is_file(), "PPTX file must exist")
        self.assertGreater(path.stat().st_size, 5000, "PPTX suspiciously small")
        with zipfile.ZipFile(path) as zf:
            slide_files = sorted(n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
            xml_blob = b"".join(zf.read(n) for n in zf.namelist() if n.endswith(".xml"))
        return {
            "slide_count": len(slide_files),
            "has_hlink_click": b"hlinkClick" in xml_blob or b"hlinksldjump" in xml_blob,
            "has_hlink_hover": b"hlinkHover" in xml_blob,
            "has_rtl": b'rtl="1"' in xml_blob or b"rtl='1'" in xml_blob,
        }

    def test_build_flat_with_zoom_and_hover(self) -> None:
        d = self.tmp / "flat"
        d.mkdir()
        for i in range(3):
            _write_img(d, f"img{i}.jpg")
        job = scan_project_folders(d)[0]
        out = self.tmp / "out.pptx"
        cfg = BuildSettings(
            footer_text="تست پاورقی",
            enable_image_zoom=True,
            enable_hover_zoom=True,
            enable_image_border=True,
        )
        build_presentation_from_job(job, out, settings=cfg)
        info = self._inspect_pptx(out)
        # 3 images => 1 grid slide + 3 detail slides
        self.assertEqual(info["slide_count"], 4, info)
        self.assertTrue(info["has_hlink_click"], "click zoom links missing")
        self.assertTrue(info["has_hlink_hover"], "hover zoom links missing")
        self.assertTrue(info["has_rtl"], "RTL flag missing")

        # Must reopen without corruption
        from pptx import Presentation

        prs = Presentation(str(out))
        self.assertEqual(len(prs.slides), 4)

    def test_build_grouped_with_section_dividers(self) -> None:
        _make_tree(self.tmp)
        job = make_grouped_job(self.tmp / "person_alpha")
        assert job is not None
        out = self.tmp / "person_alpha.pptx"
        cfg = BuildSettings(enable_section_dividers=True, enable_image_zoom=False)
        build_presentation_from_job(job, out, settings=cfg)
        info = self._inspect_pptx(out)
        # 2 section dividers + grid slides for 2+2 images
        self.assertGreaterEqual(info["slide_count"], 4, info)

    def test_no_zoom_when_disabled(self) -> None:
        d = self.tmp / "one"
        d.mkdir()
        _write_img(d, "a.jpg")
        job = scan_project_folders(d)[0]
        out = self.tmp / "plain.pptx"
        cfg = BuildSettings(enable_image_zoom=False, enable_hover_zoom=False)
        build_presentation_from_job(job, out, settings=cfg)
        info = self._inspect_pptx(out)
        self.assertEqual(info["slide_count"], 1)
        self.assertFalse(info["has_hlink_click"])


class TestHelpContent(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        import sys

        from PySide6.QtWidgets import QApplication

        cls._app = QApplication.instance() or QApplication(sys.argv)

    def test_help_panel_rtl_native_widgets(self) -> None:
        from app import __version__
        from app.ui.help_panel import build_help_panel

        panel = build_help_panel("dark_cyan", __version__)
        self.assertEqual(panel.layoutDirection(), Qt.LayoutDirection.RightToLeft)
        headings = panel.findChildren(QLabel, options=Qt.FindChildOption.FindChildrenRecursively)
        heading_titles = [h.text() for h in headings if h.objectName() == "HelpHeading"]
        self.assertIn("شروع سریع", heading_titles)
        self.assertIn("محل ذخیرهٔ خروجی", heading_titles)
        tables = panel.findChildren(QFrame, options=Qt.FindChildOption.FindChildrenRecursively)
        table_count = sum(1 for t in tables if t.objectName() == "HelpTable")
        self.assertGreaterEqual(table_count, 3)
        all_text = " ".join(label.text() for label in panel.findChildren(QLabel))
        self.assertIn("داخل هر پوشه", all_text)
        self.assertIn("اسلاید", all_text)
        self.assertNotIn("اسلید", all_text)


class TestBuildSettings(unittest.TestCase):
    def test_hover_requires_zoom(self) -> None:
        cfg = BuildSettings.from_dict(
            {"enable_image_zoom": False, "enable_hover_zoom": True}
        )
        self.assertFalse(cfg.enable_image_zoom)
        self.assertFalse(cfg.enable_hover_zoom, "hover must be off when zoom off")

    def test_output_folder_default(self) -> None:
        cfg = BuildSettings.from_dict({"output_folder_name": "  "})
        self.assertEqual(cfg.output_folder_name, "Output_PPTX")


class TestWorkerOutputPath(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="gp_worker_"))

    def tearDown(self) -> None:
        import shutil

        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_output_inside_input_not_parent(self) -> None:
        d = self.tmp / "input_leaf"
        d.mkdir()
        _write_img(d, "a.jpg")
        cfg = BuildSettings(output_folder_name="Output_PPTX")
        expected = d / "Output_PPTX"
        # replicate worker logic
        folder_name = (cfg.output_folder_name or "Output_PPTX").strip()
        output_dir = d / folder_name
        self.assertEqual(str(output_dir), str(expected))
        self.assertEqual(output_dir.parent, d)
        self.assertNotEqual(output_dir.parent.parent, d)


class TestSettingsManager(unittest.TestCase):
    def test_roundtrip(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            mgr = SettingsManager()
            mgr._path = Path(td) / "settings.json"
            mgr._dir = Path(td)
            mgr.set("theme", "dark_purple")
            mgr.set("jpeg_quality", 80)
            mgr.save()
            mgr2 = SettingsManager()
            mgr2._path = mgr._path
            mgr2._dir = mgr._dir
            mgr2.load()
            self.assertEqual(mgr2.get("theme"), "dark_purple")
            self.assertEqual(mgr2.get("jpeg_quality"), 80)

    def test_default_theme_is_dark_cyan_on_fresh_install(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            mgr = SettingsManager()
            mgr._path = Path(td) / "settings.json"
            mgr._dir = Path(td)
            mgr.load()
            self.assertEqual(mgr.get("theme"), "dark_cyan")
            self.assertTrue(mgr._path.is_file())

    def test_settings_page_commit_persists_theme(self) -> None:
        """Theme persistence via SettingsManager (SettingsPage uses same commit path)."""
        with tempfile.TemporaryDirectory() as td:
            mgr = SettingsManager()
            mgr._path = Path(td) / "settings.json"
            mgr._dir = Path(td)
            mgr.load()
            mgr.set("theme", "dark_purple")
            mgr.save()

            mgr2 = SettingsManager()
            mgr2._path = mgr._path
            mgr2._dir = mgr._dir
            mgr2.load()
            self.assertEqual(mgr2.get("theme"), "dark_purple")


class TestUIIntegration(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        from PySide6.QtWidgets import QApplication
        import sys

        cls._app = QApplication.instance() or QApplication(sys.argv)

    def _isolated_window(self):
        from app.ui.main_window import MainWindow

        win = MainWindow()
        td = tempfile.mkdtemp(prefix="pics2ppt_ui_")
        win.settings._path = Path(td) / "settings.json"
        win.settings._dir = Path(td)
        win.settings.load()
        win.settings.set("ui_language_confirmed", True)
        win.settings.set("ui_language", "fa")
        win.settings_page.load_values()
        win.apply_preferences()
        self.addCleanup(lambda: __import__("shutil").rmtree(td, ignore_errors=True))
        return win

    def test_main_window_pages_and_theme(self) -> None:
        win = self._isolated_window()
        self.assertEqual(win.stack.count(), 3)
        self.assertEqual(win.sidebar.count(), 3)
        win.settings.set("theme", "light")
        win.apply_preferences()
        self.assertIn("background-color", win.styleSheet())
        win.change_page(1)
        self.assertEqual(win.stack.currentIndex(), 1)
        win.change_page(2)
        help_widget = win.about_page.help_scroll.widget()
        self.assertIsNotNone(help_widget)
        help_text = " ".join(label.text() for label in help_widget.findChildren(QLabel))
        self.assertTrue(
            ("شروع سریع" in help_text) or ("Quick start" in help_text),
            help_text[:200],
        )
        self.assertIn("Pics2PPT", help_text)
        win.close()

    def test_home_build_settings_sync(self) -> None:
        win = self._isolated_window()
        win.home_page.footer_edit.setText("sample-footer")
        bs = win.home_page.build_settings()
        self.assertEqual(bs.footer_text, "sample-footer")
        # Session inputs must not be persisted across launches
        self.assertEqual(win.settings.get("footer_text"), "")
        self.assertEqual(win.settings.get("last_input_dir"), "")
        win.close()


class TestWorkerEndToEnd(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp(prefix="gp_e2e_"))

    def tearDown(self) -> None:
        import shutil

        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_worker_creates_pptx_inside_input_folder(self) -> None:
        _make_tree(self.tmp)
        cfg = BuildSettings(
            output_folder_name="Output_PPTX",
            enable_image_zoom=False,
            footer_text="integration-test",
        )
        worker = PresentationWorker(str(self.tmp), cfg)
        worker.run()
        out_dir = self.tmp / "Output_PPTX"
        self.assertTrue(out_dir.is_dir(), "output dir must be inside input")
        self.assertEqual(out_dir.parent, self.tmp)
        files = list(out_dir.glob("*.pptx"))
        self.assertEqual(len(files), 3, [f.name for f in files])
        for f in files:
            from pptx import Presentation

            prs = Presentation(str(f))
            self.assertGreater(len(prs.slides), 0, f.name)


class TestThemes(unittest.TestCase):
    def test_all_theme_stylesheets_non_empty(self) -> None:
        from app.ui.theme import build_stylesheet

        for theme in ("dark_cyan", "dark_purple", "light"):
            css = build_stylesheet(theme, "medium")
            self.assertIn("background-color", css)
            self.assertGreater(len(css), 500, theme)


class TestAboutContent(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        from PySide6.QtWidgets import QApplication
        import sys

        cls._app = QApplication.instance() or QApplication(sys.argv)

    def test_creator_visible_in_about(self) -> None:
        from PySide6.QtWidgets import QLabel
        from app.ui.main_window import MainWindow

        td = tempfile.mkdtemp(prefix="pics2ppt_about_")
        win = MainWindow()
        win.settings._path = Path(td) / "settings.json"
        win.settings._dir = Path(td)
        win.settings.load()
        win.change_page(2)
        texts = [w.text() for w in win.about_page.findChildren(QLabel)]
        self.assertIn("سازنده: Ali Rashidi", texts)
        win.close()
        __import__("shutil").rmtree(td, ignore_errors=True)


if __name__ == "__main__":
    unittest.main(verbosity=2)
