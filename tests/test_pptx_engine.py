"""Tests for Hybrid Smart PPTX engine (Phase 0–1)."""

from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.core.models import BuildSettings, PptxOutputSettings
from app.core.pptx.engine import BuildPath, HybridEngine, resolve_template_file
from app.core.pptx.template_analyzer import analyze_template, format_analysis_report
from app.core.pptx.template_fill import fill_presentation_tokens, replace_tokens_in_text
from app.core.pptx.template_loader import (
    bundled_template_if_available,
    default_template_path,
    is_template_file,
    validate_template_zip,
)
from app.core.scanner import scan_project_folders
from app.core.pptx import build_presentation_from_job
from app.services.settings import SETTINGS_VERSION, SettingsManager


def _write_img(folder: Path, name: str) -> None:
    from PIL import Image

    Image.new("RGB", (40, 30), color=(100, 120, 140)).save(folder / name)


class TestHybridEngine(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp())
        self.engine = HybridEngine()

    def test_code_path_by_default(self) -> None:
        d = self.tmp / "flat"
        d.mkdir()
        _write_img(d, "a.jpg")
        job = scan_project_folders(d)[0]
        cfg = PptxOutputSettings(output_mode="code")
        result = self.engine.build(job, self.tmp / "out.pptx", settings=cfg)
        self.assertEqual(result.path_used, BuildPath.CODE)
        self.assertTrue(result.output_path.is_file())

    def test_template_mode_requires_file(self) -> None:
        d = self.tmp / "one"
        d.mkdir()
        _write_img(d, "a.jpg")
        job = scan_project_folders(d)[0]
        cfg = PptxOutputSettings(output_mode="template", template_path=None)
        with patch("app.core.pptx.engine.bundled_template_if_available", return_value=None):
            with self.assertRaises(FileNotFoundError):
                self.engine.build(job, self.tmp / "out.pptx", settings=cfg)

    def test_auto_without_any_template_uses_code(self) -> None:
        with patch("app.core.pptx.engine.bundled_template_if_available", return_value=None):
            self.assertEqual(
                self.engine.resolve_path(PptxOutputSettings(output_mode="auto")),
                BuildPath.CODE,
            )

    def test_auto_with_bundled_uses_template(self) -> None:
        bundled = bundled_template_if_available()
        if bundled is None:
            self.skipTest("bundled template missing")
        self.assertEqual(
            self.engine.resolve_path(PptxOutputSettings(output_mode="auto")),
            BuildPath.TEMPLATE,
        )

    def test_template_path_builds_with_cover_tokens(self) -> None:
        bundled = bundled_template_if_available()
        if bundled is None:
            self.skipTest("bundled template missing")
        d = self.tmp / "flat"
        d.mkdir()
        _write_img(d, "a.jpg")
        job = scan_project_folders(d)[0]
        cfg = PptxOutputSettings(
            output_mode="template",
            template_path=bundled,
            footer_text="Footer X",
            enable_image_zoom=False,
        )
        result = self.engine.build(job, self.tmp / "tpl_out.pptx", settings=cfg)
        self.assertEqual(result.path_used, BuildPath.TEMPLATE)
        with zipfile.ZipFile(result.output_path) as zf:
            xml_blob = b"".join(zf.read(n) for n in zf.namelist() if n.endswith(".xml"))
        self.assertIn(b"Footer X", xml_blob)
        self.assertNotIn(b"{{title}}", xml_blob)
        self.assertNotIn(b"{{footer}}", xml_blob)


class TestRunSafeFill(unittest.TestCase):
    def test_replace_tokens_simple(self) -> None:
        self.assertEqual(replace_tokens_in_text("Hi {{title}}!", {"title": "A"}), "Hi A!")

    def test_split_runs_across_token(self) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = "{{"
        r1.font.size = Pt(14)
        r2 = p.add_run()
        r2.text = "title"
        r2.font.size = Pt(14)
        r3 = p.add_run()
        r3.text = "}}"
        r3.font.size = Pt(14)
        changed = fill_presentation_tokens(prs, {"title": "Report"})
        self.assertGreaterEqual(changed, 1)
        self.assertEqual("".join(run.text for run in p.runs), "Report")


class TestTemplateSecurity(unittest.TestCase):
    def test_bundled_passes_validation(self) -> None:
        path = default_template_path()
        if not path.is_file():
            self.skipTest("bundled missing")
        validate_template_zip(path)

    def test_reject_path_traversal_zip(self) -> None:
        evil = Path(tempfile.mkdtemp()) / "evil.pptx"
        with zipfile.ZipFile(evil, "w") as zf:
            zf.writestr("[Content_Types].xml", "<Types/>")
            zf.writestr("../escape.txt", "x")
        with self.assertRaises(ValueError):
            validate_template_zip(evil)


class TestTemplateAnalyzer(unittest.TestCase):
    def test_analyze_bundled(self) -> None:
        path = bundled_template_if_available()
        if path is None:
            self.skipTest("bundled missing")
        rows = analyze_template(path)
        self.assertGreaterEqual(len(rows), 1)
        report = format_analysis_report(rows)
        self.assertIn("layout", report.lower())


class TestPptxOutputSettings(unittest.TestCase):
    def test_from_dict_new_fields(self) -> None:
        cfg = PptxOutputSettings.from_dict({
            "output_mode": "code",
            "template_path": "",
            "slide_size_preset": "standard_4_3",
            "slide_width_inches": 10.0,
            "slide_height_inches": 7.5,
            "title_font_size": 24,
            "image_fit": "fill",
        })
        self.assertEqual(cfg.output_mode, "code")
        self.assertEqual(cfg.slide_size_preset, "standard_4_3")
        self.assertEqual(cfg.title_font_size, 24)
        self.assertEqual(cfg.image_fit, "fill")

    def test_invalid_output_mode_defaults_auto(self) -> None:
        cfg = PptxOutputSettings.from_dict({"output_mode": "invalid"})
        self.assertEqual(cfg.output_mode, "auto")


class TestFilenameXmlSafety(unittest.TestCase):
    """G28 — ampersand in filename must not break PPTX XML."""

    def setUp(self) -> None:
        self.tmp = Path(tempfile.mkdtemp())

    def test_ampersand_filename_builds_valid_pptx(self) -> None:
        d = self.tmp / "flat"
        d.mkdir()
        _write_img(d, "report&summary.jpg")
        job = scan_project_folders(d)[0]
        out = self.tmp / "out.pptx"
        build_presentation_from_job(job, out, settings=BuildSettings(caption_from_filename=True))
        self.assertTrue(out.is_file())
        with zipfile.ZipFile(out) as zf:
            xml_blob = b"".join(zf.read(n) for n in zf.namelist() if n.endswith(".xml"))
        self.assertIn(b"report", xml_blob)
        self.assertNotIn(b"&summary", xml_blob)


class TestSettingsV6Migration(unittest.TestCase):
    def test_v5_profile_migrates_pptx_keys(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            base = Path(td)
            settings_path = base / "settings.json"
            settings_path.write_text(
                '{"settings_version": 5, "theme": "dark_cyan", "ui_language_confirmed": true}',
                encoding="utf-8",
            )
            mgr = SettingsManager.__new__(SettingsManager)
            mgr._dir = base
            mgr._legacy_dirs = []
            mgr._path = settings_path
            mgr._data = {}
            mgr._fresh_install = False
            mgr.load()
            self.assertEqual(mgr.get("settings_version"), SETTINGS_VERSION)
            self.assertEqual(mgr.get("output_mode"), "auto")
            self.assertEqual(mgr.get("slide_size_preset"), "widescreen_16_9")
            self.assertEqual(mgr.get("image_fit"), "fit")


class TestTemplateLoader(unittest.TestCase):
    def test_is_template_file(self) -> None:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = Path(f.name)
        try:
            self.assertTrue(is_template_file(path))
            self.assertFalse(is_template_file(path.with_suffix(".txt")))
        finally:
            path.unlink(missing_ok=True)

    def test_resolve_template_prefers_user(self) -> None:
        bundled = bundled_template_if_available()
        if bundled is None:
            self.skipTest("bundled missing")
        cfg = PptxOutputSettings(template_path=bundled)
        self.assertEqual(resolve_template_file(cfg), bundled)


if __name__ == "__main__":
    unittest.main()
