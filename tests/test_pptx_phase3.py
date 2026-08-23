"""Phase 3 — presets, colors, index slide, template import."""

from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation

from app.core.models import PptxOutputSettings
from app.core.pptx.engine import HybridEngine
from app.core.pptx.presets import (
    apply_preset_to_mapping,
    list_builtin_presets,
    preset_settings,
    save_user_preset,
    resolve_preset_settings,
)
from app.core.pptx.template_import import import_template, layout_wizard_report
from app.core.pptx.template_loader import bundled_template_if_available
from app.core.pptx.themes import accent_color, normalize_hex, rgb_from_hex
from app.core.scanner import scan_project_folders
from pptx.dml.color import RGBColor


def _write_img(folder: Path, name: str) -> Path:
    path = folder / name
    Image.new("RGB", (40, 30), color=(80, 100, 120)).save(path)
    return path


class TestPresets(unittest.TestCase):
    def test_builtin_ids(self) -> None:
        ids = list_builtin_presets()
        self.assertEqual(set(ids), {"report", "minimal", "print", "brand"})

    def test_apply_minimal_preset(self) -> None:
        base = {"jpeg_quality": 75, "enable_image_zoom": True, "theme": "dark_cyan"}
        merged = apply_preset_to_mapping(base, "minimal")
        self.assertFalse(merged["enable_image_zoom"])
        self.assertEqual(merged["active_preset"], "minimal")
        self.assertEqual(merged["theme"], "dark_cyan")

    def test_user_preset_roundtrip(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            base = Path(td)
            save_user_preset("MyCustom", {"jpeg_quality": 88, "enable_index_slide": True}, base_dir=base)
            loaded = resolve_preset_settings("MyCustom", base_dir=base)
            self.assertEqual(loaded["jpeg_quality"], 88)
            self.assertTrue(loaded["enable_index_slide"])


class TestThemes(unittest.TestCase):
    def test_normalize_hex(self) -> None:
        self.assertEqual(normalize_hex("#0f3d2e", "000000"), "0F3D2E")
        self.assertEqual(normalize_hex("bad", "FFFFFF"), "FFFFFF")

    def test_accent_from_settings(self) -> None:
        cfg = PptxOutputSettings(color_accent="112233")
        self.assertEqual(accent_color(cfg), RGBColor(0x11, 0x22, 0x33))


class TestIndexSlide(unittest.TestCase):
    def test_index_slide_adds_table(self) -> None:
        tmp = Path(tempfile.mkdtemp())
        d = tmp / "flat"
        d.mkdir()
        _write_img(d, "a.jpg")
        _write_img(d, "b.jpg")
        job = scan_project_folders(d)[0]
        out = tmp / "idx.pptx"
        cfg = PptxOutputSettings(
            output_mode="code",
            enable_index_slide=True,
            enable_image_zoom=False,
            write_build_report=False,
        )
        HybridEngine().build(job, out, settings=cfg)
        prs = Presentation(str(out))
        # cover/index + at least one grid
        self.assertGreaterEqual(len(prs.slides), 2)
        has_table = any(shape.has_table for shape in prs.slides[0].shapes)
        self.assertTrue(has_table)


class TestTemplateImport(unittest.TestCase):
    def test_import_bundled(self) -> None:
        bundled = bundled_template_if_available()
        if bundled is None:
            self.skipTest("bundled missing")
        with tempfile.TemporaryDirectory() as td:
            dest = import_template(bundled, base_dir=Path(td))
            self.assertTrue(dest.is_file())
            report = layout_wizard_report(dest)
            self.assertIn("layout_index_grid", report)


class TestBrandColorsInOutput(unittest.TestCase):
    def test_custom_accent_survives_build(self) -> None:
        tmp = Path(tempfile.mkdtemp())
        d = tmp / "flat"
        d.mkdir()
        _write_img(d, "a.jpg")
        job = scan_project_folders(d)[0]
        # Force section divider path with a single group won't add divider;
        # just ensure build with custom colors does not crash.
        out = tmp / "colors.pptx"
        cfg = PptxOutputSettings(
            output_mode="code",
            color_accent="AABBCC",
            color_title="112233",
            enable_image_zoom=False,
            write_build_report=False,
        )
        HybridEngine().build(job, out, settings=cfg)
        self.assertTrue(out.is_file())
        with zipfile.ZipFile(out) as zf:
            xml = b"".join(zf.read(n) for n in zf.namelist() if n.endswith(".xml"))
        self.assertIn(b"112233", xml)


if __name__ == "__main__":
    unittest.main()
